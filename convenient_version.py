import os
import re
import sys
import openai

import docx
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
import PyPDF2
import pptx
import openpyxl

import tkinter as tk
from tkinter import messagebox
import io

# =========== 设置你的 openai.api_base 和 openai.api_key ===========
openai.api_base = "将此部分替换为URL"
openai.api_key = "将此部分替换为API"

def read_pdf(file_path):
    """读取 PDF 文件内容并返回纯文本。"""
    text_content = []
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                if page:
                    txt = page.extract_text()
                    if txt:
                        text_content.append(txt)
    except Exception as e:
        print(f"读取 PDF 文件出错: {file_path}, 错误信息: {e}")
    return "\n".join(text_content)

def read_docx(file_path):
    """读取 DOCX 文件内容并返回纯文本。"""
    text_content = []
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
    except Exception as e:
        print(f"读取 DOCX 文件出错: {file_path}, 错误信息: {e}")
    return "\n".join(text_content)

def read_pptx(file_path):
    """读取 PPTX 文件内容并返回纯文本。"""
    text_content = []
    try:
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text_content.append("\n".join(slide_text))
    except Exception as e:
        print(f"读取 PPTX 文件出错: {file_path}, 错误信息: {e}")
    return "\n".join(text_content)

def read_xlsx(file_path):
    """读取 XLSX 文件内容并返回纯文本。"""
    text_content = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_str.strip():
                    text_content.append(row_str)
    except Exception as e:
        print(f"读取 Excel 文件出错: {file_path}, 错误信息: {e}")
    return "\n".join(text_content)

def clean_text(text):
    """可选：去除零宽空格等特殊字符，减小字体混乱风险。"""
    text = re.sub(r'[\u200B-\u200F\uFEFF]', '', text)
    text = re.sub(r'\n\s*\n+', '\n', text)
    return text

def call_openai_api(prompt_text):
    """
    在这里：若 OpenAI 调用失败或返回空，则抛出异常，给外层捕获。
    这样就不会继续执行到报告生成，也不会打印“已成功生成报告”。
    """
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant who summarizes documents for the user."},
                {"role": "user", "content": prompt_text}
            ],
            temperature=0.7,
            max_tokens=1000
        )
    except Exception as e:
        # 遇到例如 "API Key不对"、"网络超时"、"版本不支持" 等
        print(f"调用 openai.ChatCompletion.create 出错: {e}")
        # 关键：抛出异常
        raise RuntimeError(f"调用OpenAI接口失败: {e}")

    # 成功请求，但可能choices为空
    if not response.choices:
        raise RuntimeError("OpenAI返回空结果，无法生成总结")

    # 正常情况下拿到content
    content = response.choices[0].message["content"].strip()
    if not content:
        raise RuntimeError("OpenAI返回的总结内容为空")

    return content

def generate_word_report(summary_text, output_file="综合总结与工作报告.docx"):
    """
    将给定的 summary_text 写入 Word 文档，并统一字体和字号。
    若写入出现问题，也抛出异常而不是静默返回。
    """
    if not summary_text.strip():
        raise RuntimeError("要生成的总结文本为空，无法生成报告")

    try:
        document = Document()

        style_normal = document.styles['Normal']
        style_normal.font.name = 'Microsoft YaHei'
        style_normal._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        style_normal.font.size = Pt(12)

        heading_style = document.styles['Heading 1']
        heading_style.font.name = 'Microsoft YaHei'
        heading_style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        heading_style.font.size = Pt(16)

        document.add_heading("综合总结与工作报告", level=1)

        paragraphs = summary_text.split('\n')
        for paragraph_text in paragraphs:
            paragraph_text = paragraph_text.strip()
            if not paragraph_text:
                continue
            p = document.add_paragraph(style='Normal')
            run = p.add_run(paragraph_text)
            run.font.name = 'Microsoft YaHei'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
            run.font.size = Pt(12)

        document.save(output_file)
        # 真正写入成功后，才打印这句【唯一成功标志】:
        print(f"已成功生成报告: {output_file}")
    except Exception as e:
        print(f"生成 Word 报告出错: {e}")
        raise RuntimeError(f"生成 Word 报告出错: {e}")

def main():
    """主函数：扫描当前目录文档 -> 调用API总结 -> 生成报告。"""
    current_dir = os.getcwd()
    files = os.listdir(current_dir)

    all_texts = []
    supported_ext = [".pdf", ".docx", ".pptx", ".xlsx"]
    for f in files:
        file_path = os.path.join(current_dir, f)
        if not os.path.isfile(file_path):
            continue

        _, ext = os.path.splitext(f)
        ext = ext.lower()

        file_content = ""
        if ext == ".pdf":
            file_content = read_pdf(file_path)
        elif ext == ".docx":
            file_content = read_docx(file_path)
        elif ext == ".pptx":
            file_content = read_pptx(file_path)
        elif ext == ".xlsx":
            file_content = read_xlsx(file_path)

        # 清洗文本
        file_content = clean_text(file_content)

        if file_content.strip():
            all_texts.append(f"【文件名: {f}】\n{file_content}\n")

    if not all_texts:
        # 没找到文件 or 文件空
        raise RuntimeError("当前目录下未找到可读取的文档，或文档内容为空。")

    all_text_combined = "\n\n".join(all_texts)
    prompt_text = (
        "以下是我需要你帮忙总结的多个文档内容，请整理出一个综合摘要和工作报告：\n\n"
        + all_text_combined
    )
    prompt_text = clean_text(prompt_text)

    summary = call_openai_api(prompt_text)
    generate_word_report(summary, output_file="综合总结与工作报告.docx")

def run_with_popup():
    """
    执行 main() 并在弹窗中显示结果：只有当 main() 全部成功且打印了
    “已成功生成报告: 综合总结与工作报告.docx”时才视为成功。
    任何异常或不打印这行，都视为失败。
    """
    root = tk.Tk()
    root.withdraw()

    import io
    buffer = io.StringIO()
    old_stdout = sys.stdout
    sys.stdout = buffer

    try:
        main()
        logs = buffer.getvalue()
        # 判断是否包含成功标志
        success_marker = "已成功生成报告: 综合总结与工作报告.docx"
        if success_marker in logs:
            messagebox.showinfo("执行结果", f"程序已运行成功！\n\n---输出信息---\n{logs}")
        else:
            messagebox.showerror("执行结果", f"程序运行不对啊哥们，仔细检查检查啊！\n\n---输出信息---\n{logs}")

    except Exception as e:
        # 任何抛出的异常都到这里
        logs = buffer.getvalue()
        messagebox.showerror("执行结果", f"程序运行不对啊哥们，仔细检查检查啊！\n\n---错误---\n{e}\n\n---输出信息---\n{logs}")

    finally:
        sys.stdout = old_stdout
        root.destroy()

if __name__ == "__main__":
    run_with_popup()

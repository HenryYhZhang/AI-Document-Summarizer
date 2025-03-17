import os
import re
import sys
import time
import openai
import winreg
import tkinter as tk
from tkinter import messagebox, ttk
import docx
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
import PyPDF2
import pptx
import openpyxl
import io

#######################################
# 注册表存取 API 配置部分
#######################################
REG_PATH = r"Software\MyAPIConfig"

def save_to_registry(api_base, api_key):
    """将 API 信息存入注册表 HKEY_CURRENT_USER\Software\MyAPIConfig"""
    try:
        key = winreg.CreateKey(winreg.HKEY_CURRENT_USER, REG_PATH)
        winreg.SetValueEx(key, "api_base", 0, winreg.REG_SZ, api_base)
        winreg.SetValueEx(key, "api_key", 0, winreg.REG_SZ, api_key)
        winreg.CloseKey(key)
    except Exception as e:
        print(f"[DEBUG] 写入注册表失败: {e}")

def load_from_registry():
    """从注册表读取 API 信息 (api_base, api_key)"""
    try:
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, REG_PATH)
        api_base, _ = winreg.QueryValueEx(key, "api_base")
        api_key, _ = winreg.QueryValueEx(key, "api_key")
        winreg.CloseKey(key)
        return api_base, api_key
    except FileNotFoundError:
        return None, None

def ask_for_new_api_info():
    """
    弹出配置窗口让用户输入 API URL 和 API Key（两个输入框在同一窗口）。
    提示文本：
    “自己去 GitHub 上搜一下类似 free_chatgpt_api 的字样。干嘛？没有还找我伸手啊？”
    用户点击【确定】后保存到注册表；点击【取消】则退出程序。
    """
    def on_ok():
        base = url_entry.get().strip()
        key_ = key_entry.get().strip()
        if not base or not key_:
            messagebox.showerror("错误", "URL 和 API Key 不能为空！")
            return
        save_to_registry(base, key_)
        config_win.api_info = (base, key_)
        config_win.destroy()

    def on_cancel():
        config_win.api_info = None
        config_win.destroy()

    config_win = tk.Toplevel()
    config_win.title("API 配置")
    config_win.geometry("400x240")
    config_win.resizable(False, False)

    tk.Label(config_win, text="请输入 API URL 和 API Key", font=("Arial", 12)).pack(pady=5)
    tk.Label(config_win, text="自己去 GitHub 上搜一下类似 free_chatgpt_api 的字样。\n干嘛？没有还找我伸手啊？",
             font=("Arial", 10), fg="gray").pack(pady=5)
    tk.Label(config_win, text="URL:", font=("Arial", 10)).pack()
    url_entry = tk.Entry(config_win, width=50)
    url_entry.pack(pady=5)
    tk.Label(config_win, text="API:", font=("Arial", 10)).pack()
    key_entry = tk.Entry(config_win, width=50, show="*")
    key_entry.pack(pady=5)
    btn_frame = tk.Frame(config_win)
    btn_frame.pack(pady=10)
    tk.Button(btn_frame, text="确定", width=12, height=2, command=on_ok, font=("Arial", 10)).pack(side="left", padx=10)
    tk.Button(btn_frame, text="取消", width=12, height=2, command=on_cancel, font=("Arial", 10)).pack(side="left", padx=10)

    config_win.api_info = None
    config_win.wait_window()
    if config_win.api_info is None:
        sys.exit(0)
    return config_win.api_info

def load_or_ask_api_info():
    """
    尝试从注册表读取 API 信息；若不存在，则弹出输入窗口。
    """
    api_base, api_key = load_from_registry()
    if api_base and api_key:
        return api_base, api_key
    return ask_for_new_api_info()

#######################################
# 获取程序所在目录（兼容 exe 和脚本）
#######################################
def get_base_dir():
    """
    若是 PyInstaller 打包后的 exe 运行，则用 sys.executable；
    否则用 __file__。
    """
    if getattr(sys, 'frozen', False):  # 表示已打包为 exe
        return os.path.dirname(sys.executable)
    else:
        return os.path.dirname(os.path.abspath(__file__))

#######################################
# 欢迎窗口（始终显示一次，无需传参）
#######################################
def show_welcome_window():
    """
    欢迎窗口：显示“欢迎使用”四个字（居中、黑体加粗、字号大），
    下方依次显示：
      “确保此程序与所需总结文件在同一目录再运行”
      “初次使用请先配置 URL/API”
    并在窗口左下角显示版本号“V2.0”，右下角显示“作者：HenryYhZhang”。
    包含两个按钮：
      - “开始总结”：点击后关闭欢迎窗口，程序继续；
      - “配置 URL/API”：点击后弹出配置窗口更新 API 信息，但欢迎窗口仍保持显示。
    如果用户关闭欢迎窗口，则终止程序。
    """
    def on_start():
        win_welcome.result = "start"
        win_welcome.destroy()

    def on_config():
        new_info = ask_for_new_api_info()
        if new_info:
            openai.api_base, openai.api_key = new_info
        # 不关闭欢迎窗口，等待用户点击“开始总结”

    def on_close():
        win_welcome.destroy()
        sys.exit(0)

    win_welcome = tk.Toplevel()
    win_welcome.title("欢迎")
    win_welcome.geometry("400x250")
    win_welcome.resizable(False, False)
    win_welcome.protocol("WM_DELETE_WINDOW", on_close)

    # 标题：欢迎使用（居中、黑体加粗，字号大）
    tk.Label(win_welcome, text="欢迎使用", font=("SimHei", 24, "bold")).pack(pady=10)

    # 新增说明：确保此程序与所需总结文件在同一目录再运行
    tk.Label(win_welcome, text="确保此程序与所需总结文件在同一目录再运行", font=("SimHei", 10), fg="gray").pack(pady=5)
    # 说明：初次使用请先配置 URL/API
    tk.Label(win_welcome, text="初次使用请先配置 URL/API", font=("SimHei", 10), fg="gray").pack(pady=5)

    btn_frame = tk.Frame(win_welcome)
    btn_frame.pack(pady=15)
    tk.Button(btn_frame, text="开始总结", width=12, height=2, command=on_start, font=("Arial", 10)).pack(side="left", padx=10)
    tk.Button(btn_frame, text="配置 URL/API", width=12, height=2, command=on_config, font=("Arial", 10)).pack(side="left", padx=10)

    footer = tk.Frame(win_welcome)
    footer.pack(side="bottom", fill="x", padx=10, pady=5)
    tk.Label(footer, text="V2.0", font=("Arial", 8)).pack(side="left")
    tk.Label(footer, text="作者：HenryYhZhang", font=("Arial", 8)).pack(side="right")

    win_welcome.wait_window()
    if not hasattr(win_welcome, "result") or win_welcome.result != "start":
        sys.exit(0)

#######################################
# 全局初始化 API 配置（隐藏主窗口）
#######################################
root_main = tk.Tk()
root_main.withdraw()
api_base, api_key = load_or_ask_api_info()
if not api_base or not api_key:
    sys.exit(0)

# 配置 OpenAI
openai.api_base = api_base
openai.api_key = api_key

#######################################
# 文件读取与清洗部分
#######################################
def read_pdf(file_path):
    texts = []
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                txt = page.extract_text()
                if txt:
                    texts.append(txt)
    except Exception as e:
        print(f"[DEBUG] 读取 PDF 失败: {file_path}, {e}")
    return "\n".join(texts)

def read_docx(file_path):
    texts = []
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
    except Exception as e:
        print(f"[DEBUG] 读取 DOCX 失败: {file_path}, {e}")
    return "\n".join(texts)

def read_pptx(file_path):
    texts = []
    try:
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                texts.append("\n".join(slide_text))
    except Exception as e:
        print(f"[DEBUG] 读取 PPTX 失败: {file_path}, {e}")
    return "\n".join(texts)

def read_xlsx(file_path):
    texts = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_str.strip():
                    texts.append(row_str)
    except Exception as e:
        print(f"[DEBUG] 读取 XLSX 失败: {file_path}, {e}")
    return "\n".join(texts)

def clean_text(text):
    text = re.sub(r'[\u200B-\u200F\uFEFF]', '', text)
    text = re.sub(r'\n\s*\n+', '\n', text)
    return text

def collect_files_content():
    """
    在这里使用 get_base_dir() 获取 exe 或脚本所在目录，扫描其下的 PDF、DOCX、PPTX、XLSX 文件，
    读取文本并合并，同时打印调试信息。
    """
    base_dir = get_base_dir()
    print(f"[DEBUG] 程序所在目录: {base_dir}")

    all_texts = []
    supported_ext = [".pdf", ".docx", ".pptx", ".xlsx"]
    for f in os.listdir(base_dir):
        file_path = os.path.join(base_dir, f)
        if not os.path.isfile(file_path):
            continue
        ext = os.path.splitext(f)[1].lower()
        if ext not in supported_ext:
            print(f"[DEBUG] 跳过不支持的文件: {f}")
            continue
        content = ""
        if ext == ".pdf":
            content = read_pdf(file_path)
        elif ext == ".docx":
            content = read_docx(file_path)
        elif ext == ".pptx":
            content = read_pptx(file_path)
        elif ext == ".xlsx":
            content = read_xlsx(file_path)
        content = clean_text(content)
        print(f"[DEBUG] 文件: {f}, 扩展名: {ext}, 文本长度: {len(content)}")
        if content.strip():
            all_texts.append(content)
    return "\n\n".join(all_texts)

#######################################
# 解决 exe 路径问题：使用 sys.executable
#######################################
def get_base_dir():
    if getattr(sys, 'frozen', False):
        # PyInstaller打包后的exe环境
        return os.path.dirname(sys.executable)
    else:
        # 普通脚本环境
        return os.path.dirname(os.path.abspath(__file__))

#######################################
# 进度条、API 调用 & 报告生成部分
#######################################
def show_loading_bar():
    """
    显示假加载进度条窗口：
      - 先加载到约 61.8%
      - 返回 (root_bar, loading_win, progress)
    """
    root_bar = tk.Tk()
    root_bar.withdraw()

    loading_win = tk.Toplevel()
    loading_win.title("正在总结")
    loading_win.geometry("400x150")
    loading_win.resizable(False, False)

    tk.Label(loading_win, text="正在总结...", font=("Arial", 12)).pack(pady=10)
    progress = ttk.Progressbar(loading_win, orient="horizontal", length=300, mode="determinate")
    progress.pack(pady=5)
    tk.Label(loading_win, text="其实这是一个假的进度条", font=("Arial", 9), fg="gray").pack(pady=5)
    loading_win.update()

    for i in range(31):
        time.sleep(0.01)
        progress["value"] = 61.8 * (i / 30.0)
        loading_win.update()
    return root_bar, loading_win, progress

def call_openai_api(prompt_text):
    """
    构造 prompt 后调用 OpenAI 接口生成综合总结。
    若调用失败，则允许用户选择重新输入 URL & API Key（最多重试 3 次）。
    """
    # 在此处可直接使用 prompt_text，因为 collect_files_content() 已在 main() 中被调用
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "你是一个文档总结助手，帮助用户生成综合总结。"},
                {"role": "user", "content": prompt_text},
            ],
            temperature=0.7,
            max_tokens=1000
        )
        if not response.choices:
            raise RuntimeError("OpenAI 返回空结果")
        content = response.choices[0].message["content"].strip()
        if not content:
            raise RuntimeError("OpenAI 返回的总结内容为空")
        return content
    except Exception as e:
        # 提示是否重新输入
        retry = messagebox.askyesno("API 出错", f"调用 OpenAI 接口失败：\n{e}\n\n是否重新输入 URL 和 API Key？")
        if retry:
            new_info = ask_for_new_api_info()
            openai.api_base, openai.api_key = new_info
            # 重新发起一次
            return call_openai_api(prompt_text)
        else:
            raise RuntimeError(f"调用 OpenAI 接口失败: {e}")

def generate_word_report(summary_text, output_file="综合总结.docx"):
    """
    将综合总结写入 Word 文档，统一使用微软雅黑 12pt，
    标题为“综合总结”，保存到exe/脚本所在目录，并打印成功标志。
    """
    if not summary_text.strip():
        raise RuntimeError("要生成的总结文本为空，无法生成报告")
    try:
        document = Document()

        style_normal = document.styles['Normal']
        style_normal.font.name = 'Microsoft YaHei'
        style_normal._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        style_normal.font.size = Pt(12)

        heading_style = document.styles['Heading 1']
        heading_style.font.name = 'Microsoft YaHei'
        heading_style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        heading_style.font.size = Pt(16)

        document.add_heading("综合总结", level=1)

        for paragraph in summary_text.split('\n'):
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            p = document.add_paragraph(paragraph, style='Normal')
            p.paragraph_format.line_spacing = 1.5

        base_dir = get_base_dir()
        output_path = os.path.join(base_dir, output_file)
        document.save(output_path)
        print(f"已成功生成报告: {output_path}")
    except Exception as e:
        raise RuntimeError(f"生成 Word 报告出错: {e}")

#######################################
# 主程序逻辑
#######################################
def main():
    """
    1. 显示假加载进度条窗口（先加载到 61.8%）
    2. 调用 collect_files_content() 获取所有文件内容
    3. 构造 prompt 并调用 OpenAI 接口生成综合总结
    4. 生成 Word 文档“综合总结.docx”
    5. 补满进度条到 100%，等待 0.5 秒后关闭加载窗口
    """
    root_bar, loading_win, progress = show_loading_bar()

    merged_text = collect_files_content()
    if not merged_text.strip():
        raise RuntimeError("当前目录下未找到可读取的文档，或文档内容为空。")

    prompt_text = "以下是我需要你帮忙总结的多个文档内容，请整理出一个综合摘要：\n\n" + merged_text
    summary = call_openai_api(prompt_text)
    generate_word_report(summary, output_file="综合总结.docx")

    base_val = progress["value"]
    for i in range(31):
        time.sleep(0.01)
        progress["value"] = base_val + (100 - base_val) * (i / 30.0)
        loading_win.update()

    time.sleep(0.5)
    loading_win.destroy()
    root_bar.destroy()

#######################################
# 结果弹窗包装
#######################################
def run_with_popup():
    # 显示欢迎窗口；如果用户关闭欢迎窗口，则程序终止
    show_welcome_window()

    # 主逻辑执行
    root = tk.Tk()
    root.withdraw()

    buffer = io.StringIO()
    old_stdout = sys.stdout
    sys.stdout = buffer

    try:
        main()
        logs = buffer.getvalue()
        success_marker = "已成功生成报告:"
        if success_marker in logs:
            messagebox.showinfo("执行结果", f"程序已运行成功！\n\n---输出信息---\n{logs}")
        else:
            messagebox.showerror("执行结果", f"程序运行不对，请仔细检查输出信息：\n\n{logs}")
    except Exception as e:
        logs = buffer.getvalue()
        messagebox.showerror("执行结果",
            f"程序运行出错，请仔细检查输出信息：\n\n---错误---\n{e}\n\n---输出信息---\n{logs}")
    finally:
        sys.stdout = old_stdout
        root.destroy()

#######################################
# 程序入口
#######################################
def run_program():
    run_with_popup()

if __name__ == "__main__":
    run_program()
